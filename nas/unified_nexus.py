import requests
import logging
import json
from flask import Flask, jsonify, request, send_file, send_from_directory, Blueprint, redirect, Response, render_template_string
from flask_cors import CORS
from dotenv import load_dotenv
# Version — defined here for direct script execution support
__version__ = '1.2.2'
RELEASE_DATE = '2026-07-16'
import psutil
import platform
import time
import os
import sys
import io
import threading
import glob
import subprocess
import shutil
from pathlib import Path
from functools import wraps

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s', 
                    stream=sys.stdout)
logger = logging.getLogger(__name__)

load_dotenv()

app = Flask(__name__)
CORS(app)
app.json.sort_keys = False

# --- CONFIG ---
import os
_ROOT_FALLBACK = os.environ.get('NAS_ROOT_DIR', '')
ROOT_DIR = _ROOT_FALLBACK if _ROOT_FALLBACK else '/home/jerry/workspace'
NAS_PASSWORD = os.environ.get('NAS_PASSWORD')
if not NAS_PASSWORD:
    import logging as _log
    _log.warning('⚠️ NAS_PASSWORD not set — using insecure fallback. Set env var NAS_PASSWORD!')
    NAS_PASSWORD = 'JERRY_NEXUS_2026'

def init_app(root, password, port):
    global ROOT_DIR, NAS_PASSWORD
    ROOT_DIR = os.path.abspath(root)
    NAS_PASSWORD = password
    logger.info(f'NAS App initialized with absolute root={ROOT_DIR}, port={port}')


# --- TRASH CONFIG ---
TRASH_DIR = os.path.join(ROOT_DIR, '.trash')
TRASH_META_FILE = os.path.join(TRASH_DIR, '.trash_meta.json')

def get_trash_meta():
    """Load trash metadata JSON."""
    if not os.path.exists(TRASH_META_FILE):
        return {'items': []}
    try:
        with open(TRASH_META_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return {'items': []}

def save_trash_meta(meta):
    """Save trash metadata JSON."""
    os.makedirs(TRASH_DIR, exist_ok=True)
    with open(TRASH_META_FILE, 'w', encoding='utf-8') as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

def resolve_path(path):
    logger.info(f'[PATH-DEBUG] Input path: "{path}", Current ROOT_DIR: "{ROOT_DIR}"')
    if not path:
        res = ROOT_DIR
    elif os.path.isabs(path) and path.startswith(ROOT_DIR):
        res = os.path.normpath(path)
    else:
        clean_path = path
        if ':' in path:
            parts = path.split(':', 1)
            if len(parts) == 2:
                clean_path = parts[1].lstrip('\\/')
        
        # Primary attempt: Join with ROOT_DIR
        res = os.path.normpath(os.path.join(ROOT_DIR, clean_path.lstrip('/')))
        
        # Smart Fallback: If file doesn't exist, try searching in the NAS app directory
        if not os.path.exists(res):
            nas_app_dir = os.path.join(ROOT_DIR, 'nas_tool', 'nas')
            fallback_res = os.path.normpath(os.path.join(nas_app_dir, clean_path.lstrip('/')))
            if os.path.exists(fallback_res):
                logger.info(f'[PATH-DEBUG] Primary failed, fallback success: {fallback_res}')
                res = fallback_res
                
    logger.info(f'[PATH-DEBUG] Resolved to: "{res}"')
    return res

# Image/Video extensions for thumbnails
IMG_EXTS = {'.jpg','.jpeg','.png','.gif','.webp'}
VID_EXTS = {'.mp4','.mov','.avi','.mkv','.webm'}
DOC_EXTS = '.doc,.docx,.xls,.xlsx,.ppt,.pptx'.split(',')

def require_auth(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        provided_pw = request.args.get('password')
        if not provided_pw and request.is_json:
            provided_pw = request.json.get('password')
            
        if provided_pw != NAS_PASSWORD:
            return jsonify({"error": "Invalid or missing password"}), 401
        return f(*args, **kwargs)
    return decorated_function

# --- NAS BLUEPRINT (Handles /nas prefix) ---
nas_bp = Blueprint('nas', __name__, url_prefix='/nas')

@nas_bp.route('/api/auth', methods=['POST'])
def verify_auth():
    data = request.json
    pw = data.get('password', '')
    if pw == NAS_PASSWORD:
        return jsonify({"success": True})
    return jsonify({"success": False, "error": "Invalid password"}), 401

@nas_bp.route('/')
def serve_nas_index():
    current_dir = os.path.dirname(os.path.abspath(__file__))
    return send_from_directory(current_dir, 'index.html')

@nas_bp.route('/api/config')
def nas_config():
    """Return current NAS root path config."""
    return jsonify({
        'root': ROOT_DIR,
        'version': __version__,
        'release_date': RELEASE_DATE
    })


@nas_bp.route('/api/sysinfo')
@require_auth
def system_info():
    """Return system info with cross-platform physical disk detection."""
    # Collect all mount candidates
    all_candidates = list(psutil.disk_partitions())

    # WSL fallback: scan /mnt for Windows drives that psutil may miss
    for mp in sorted(glob.glob('/mnt/*')):
        if os.path.isdir(mp):
            class P:
                def __init__(self, mountpoint, device, fstype, opts):
                    self.mountpoint = mountpoint
                    self.device = device
                    self.fstype = fstype
                    self.opts = opts
            all_candidates.append(P(mp, mp, 'drvfs', 'rw'))

    # Fingerprint-based dedup + system noise filter
    seen_mounts = set()
    seen_fingerprints = set()
    disks = []

    # Sort: drive letters (C, D, etc) first so they survive dedup
    order = sorted(all_candidates, key=lambda x: ('/' not in x.mountpoint and len(x.mountpoint) <= 3, x.mountpoint), reverse=True)

    for p in order:
        mount = p.mountpoint
        if mount in seen_mounts:
            continue
        seen_mounts.add(mount)

        # Skip system-virt FS types
        fstype_lower = (p.fstype or '').lower()
        if any(x in fstype_lower for x in ['tmpfs', 'squashfs', 'devtmpfs', 'loop', 'overlay']):
            continue

        # Skip system root and internal WSL paths
        if mount == '/' or 'wslg' in mount or 'distro' in mount.lower() or 'snap' in mount.lower():
            continue

        try:
            usage = psutil.disk_usage(mount)
        except (PermissionError, OSError):
            continue

        # Skip truly tiny partitions (< 10 GB) — likely system partitions
        if usage.total < 10 * 1024**3:
            continue

        # Label derivation
        label = mount.split('/')[-1] if '/' in mount else mount
        if not label:
            label = mount
        label = label.upper()

        # Fingerprint dedup
        fp = (round(usage.total / 1024**3, 2), round(usage.used / 1024**3, 2))
        if fp in seen_fingerprints:
            continue
        seen_fingerprints.add(fp)

        disks.append({
            'mount': mount, 'label': label,
            'total_gb': usage.total / 1024**3,
            'used_gb': usage.used / 1024**3,
            'percent': usage.percent
        })

    cpu_percent = psutil.cpu_percent(interval=0.1)
    cpu_count = psutil.cpu_count()
    mem = psutil.virtual_memory()
    net = psutil.net_io_counters()
    boot_time = psutil.boot_time()
    uptime_seconds = time.time() - boot_time
    uptime_human = f'{int(uptime_seconds // 86400)}d {int((uptime_seconds % 86400) // 3600)}h {int((uptime_seconds % 3600) // 60)}m'

    return jsonify({
        'cpu_percent': cpu_percent, 'cpu_count': cpu_count,
        'platform': platform.platform(), 'hostname': platform.node(),
        'memory': {
            'total_gb': mem.total / 1024**3,
            'used_gb': (mem.total - mem.available) / 1024**3,
            'percent': mem.percent
        },
        'disks': disks,
        'network': {'bytes_sent': net.bytes_sent, 'bytes_recv': net.bytes_recv},
        'uptime_human': uptime_human
    })



@nas_bp.route('/api/read')
@require_auth
def read_file():
    path = request.args.get('path', '')
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR): return jsonify({"error": "Forbidden"}), 403
    try:
        with open(full_path, 'r', encoding='utf-8') as f: return f.read()
    except Exception as e: return jsonify({"error": str(e)}), 500

@nas_bp.route('/api/save', methods=['POST'])
@require_auth
def save_file():
    data = request.json
    path = data['path']
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR): return jsonify({"error": "Forbidden"}), 403
    try:
        content = data['content']
        is_binary = data.get('binary', False)
        if is_binary:
            import base64
            with open(full_path, 'wb') as f: f.write(base64.b64decode(content))
        else:
            with open(full_path, 'w', encoding='utf-8') as f: f.write(content)
        return jsonify({"success": True})
    except Exception as e: return jsonify({"error": str(e)}), 500

@nas_bp.route('/api/files')
@require_auth
def list_files():
    path = request.args.get('path', '')
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR): return jsonify({"error": "Forbidden"}), 403
    if not os.path.isdir(full_path): return jsonify({"error": "Path is not a directory"}), 400
    try:
        files = []
        for entry in os.scandir(full_path):
            try:
                st = entry.stat()
                files.append({
                    "name": entry.name, "is_dir": entry.is_dir(),
                    "size": st.st_size if not entry.is_dir() else None,
                    "mtime": st.st_mtime
                })
            except (FileNotFoundError, OSError):
                # Skip broken symlinks or inaccessible entries
                files.append({
                    "name": entry.name, "is_dir": False,
                    "size": 0, "mtime": 0
                })
        return jsonify(files)
    except Exception as e: return jsonify({"error": str(e)}), 500

@nas_bp.route('/api/mkdir', methods=['POST'])
@require_auth
def make_dir():
    data = request.json
    path = data['path']
    full_dir = resolve_path(path)
    full_path = os.path.join(full_dir, data['name'])
    try:
        os.makedirs(full_path, exist_ok=True)
        return jsonify({"success": True})
    except Exception as e: return jsonify({"error": str(e)}), 500

@nas_bp.route('/api/upload', methods=['POST'])
@require_auth
def upload_file():
    file = request.files['file']
    path = request.form.get('path', '')
    full_dir = resolve_path(path)
    
    # Handle duplicate filenames by adding numerical suffix
    filename = file.filename
    name, ext = os.path.splitext(filename)
    full_path = os.path.join(full_dir, filename)
    
    counter = 1
    while os.path.exists(full_path):
        new_filename = f"{name}_{counter}{ext}"
        full_path = os.path.join(full_dir, new_filename)
        counter += 1
        
    try:
        file.save(full_path)
        return jsonify({"success": True, "saved_path": full_path})
    except Exception as e: return jsonify({"error": str(e)}), 500

@nas_bp.route('/api/delete', methods=['POST'])
@require_auth
def delete_item():
    path = request.json.get('path', '')
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR): return jsonify({"error": "Forbidden"}), 403
    try:
        # Move to trash instead of permanent delete
        os.makedirs(TRASH_DIR, exist_ok=True)
        
        # Get relative path for metadata (to restore later)
        rel_path = os.path.relpath(full_path, ROOT_DIR)
        
        # Generate unique trash name with timestamp
        ts = int(time.time() * 1000)
        base_name = os.path.basename(full_path)
        trash_name = f"{base_name}_{ts}"
        trash_path = os.path.join(TRASH_DIR, trash_name)
        
        # Move to trash
        import shutil
        shutil.move(full_path, trash_path)
        
        # Record metadata
        meta = get_trash_meta()
        meta['items'].append({
            'trash_name': trash_name,
            'original_path': rel_path,
            'is_dir': os.path.isdir(trash_path),
            'deleted_at': time.strftime('%Y-%m-%dT%H:%M:%S', time.localtime(ts / 1000)),
            'size': os.path.getsize(trash_path) if os.path.isfile(trash_path) else None
        })
        save_trash_meta(meta)
        
        return jsonify({"success": True, "trashed": trash_name})
    except Exception as e: return jsonify({"error": str(e)}), 500


# === TRASH API ===

@nas_bp.route('/api/trash/list', methods=['GET'])
@require_auth
def trash_list():
    """List all items in the trash."""
    meta = get_trash_meta()
    # Validate each item still physically exists
    valid_items = []
    for item in meta.get('items', []):
        trash_full = os.path.join(TRASH_DIR, item['trash_name'])
        if os.path.exists(trash_full):
            valid_items.append(item)
        else:
            logger.warning(f"Trash item missing from disk: {item['trash_name']}")
    
    if len(valid_items) != len(meta.get('items', [])):
        meta['items'] = valid_items
        save_trash_meta(meta)
    
    return jsonify({'items': sorted(valid_items, key=lambda x: x.get('deleted_at', ''), reverse=True)})


@nas_bp.route('/api/trash/restore', methods=['POST'])
@require_auth
def trash_restore():
    """Restore a single item from trash to its original location."""
    trash_name = request.json.get('trash_name', '')
    if not trash_name:
        return jsonify({"error": "Missing trash_name"}), 400
    
    meta = get_trash_meta()
    item = None
    for idx, it in enumerate(meta.get('items', [])):
        if it['trash_name'] == trash_name:
            item = it
            item_idx = idx
            break
    
    if not item:
        return jsonify({"error": "Item not found in trash metadata"}), 404
    
    trash_full = os.path.join(TRASH_DIR, item['trash_name'])
    if not os.path.exists(trash_full):
        meta['items'].pop(item_idx)
        save_trash_meta(meta)
        return jsonify({"error": "Item no longer exists on disk"}), 404
    
    # Compute original full path
    original_full = os.path.join(ROOT_DIR, item['original_path'])
    original_dir = os.path.dirname(original_full)
    
    # Create parent directory if it no longer exists
    os.makedirs(original_dir, exist_ok=True)
    
    # Handle filename conflict at destination
    if os.path.exists(original_full):
        base, ext = os.path.splitext(item['original_path'])
        counter = 1
        while True:
            new_name = f"{base}_restored_{counter}{ext}"
            candidate = os.path.join(ROOT_DIR, new_name)
            if not os.path.exists(candidate):
                original_full = candidate
                break
            counter += 1
        logger.info(f"Restore conflict resolved: renamed to {original_full}")
    
    import shutil
    shutil.move(trash_full, original_full)
    
    # Remove from metadata
    meta['items'].pop(item_idx)
    save_trash_meta(meta)
    
    return jsonify({"success": True, "restored_to": os.path.relpath(original_full, ROOT_DIR)})


@nas_bp.route('/api/trash/permanent-delete', methods=['POST'])
@require_auth
def trash_permanent_delete():
    """Permanently delete a single item from trash."""
    trash_name = request.json.get('trash_name', '')
    if not trash_name:
        return jsonify({"error": "Missing trash_name"}), 400
    
    trash_full = os.path.join(TRASH_DIR, trash_name)
    
    # Delete from disk
    try:
        if os.path.isdir(trash_full):
            import shutil
            shutil.rmtree(trash_full)
        elif os.path.isfile(trash_full):
            os.remove(trash_full)
    except Exception as e:
        return jsonify({"error": f"Failed to delete from disk: {str(e)}"}), 500
    
    # Remove from metadata
    meta = get_trash_meta()
    meta['items'] = [it for it in meta.get('items', []) if it['trash_name'] != trash_name]
    save_trash_meta(meta)
    
    return jsonify({"success": True})


@nas_bp.route('/api/trash/empty', methods=['POST'])
@require_auth
def trash_empty():
    """Empty entire trash — permanently delete everything."""
    meta = get_trash_meta()
    errors = []
    
    for item in meta.get('items', []):
        trash_full = os.path.join(TRASH_DIR, item['trash_name'])
        try:
            if os.path.isdir(trash_full):
                import shutil
                shutil.rmtree(trash_full)
            elif os.path.isfile(trash_full):
                os.remove(trash_full)
        except Exception as e:
            errors.append(item['trash_name'])
            logger.warning(f"Failed to delete {item['trash_name']}: {e}")
    
    # Clear metadata
    meta['items'] = []
    save_trash_meta(meta)
    
    # Also clean up the trash directory itself if empty
    try:
        if os.path.isdir(TRASH_DIR) and not os.listdir(TRASH_DIR):
            os.rmdir(TRASH_DIR)
    except Exception:
        pass
    
    if errors:
        return jsonify({"success": True, "warning": f"{len(errors)} item(s) could not be deleted"})
    return jsonify({"success": True})

@nas_bp.route('/api/rename', methods=['POST'])
@require_auth
def rename_item():
    data = request.json
    old_p = data['oldPath']
    new_p = data['newName']
    old_full = resolve_path(old_p)
    if not old_full.startswith(ROOT_DIR): return jsonify({"error": "Forbidden"}), 403
    new_full = os.path.join(os.path.dirname(old_full), new_p)
    try:
        os.rename(old_full, new_full)
        return jsonify({"success": True})
    except Exception as e: return jsonify({"error": str(e)}), 500

@nas_bp.route('/api/download')
@require_auth
def download_file():
    path = request.args.get('path', '')
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR): return jsonify({"error": "Forbidden"}), 403
    return send_from_directory(os.path.dirname(full_path), os.path.basename(full_path), as_attachment=True)


@nas_bp.route('/api/download-zip', methods=['POST'])
@require_auth
def download_zip():
    """Download multiple files/folders as a ZIP archive."""
    import zipfile
    import tempfile
    from io import BytesIO
    
    paths = request.json.get('paths', [])
    if not paths:
        return jsonify({"error": "No paths provided"}), 400
    
    # Resolve all paths and validate
    resolved = []
    for p in paths:
        fp = resolve_path(p)
        if not fp.startswith(ROOT_DIR):
            return jsonify({"error": f"Forbidden: {p}"}), 403
        if not os.path.exists(fp):
            return jsonify({"error": f"Not found: {p}"}), 404
        resolved.append((p, fp))
    
    # Create ZIP in memory
    buf = BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for orig_path, full_path in resolved:
            # Determine archive name: use the relative path from resolve
            rel_name = orig_path.lstrip('/')
            if os.path.isdir(full_path):
                # Walk directory and add all files
                for dirpath, dirnames, filenames in os.walk(full_path):
                    # Compute archive subpath
                    sub_rel = os.path.relpath(dirpath, os.path.dirname(full_path))
                    for fn in filenames:
                        file_full = os.path.join(dirpath, fn)
                        arcname = os.path.join(os.path.basename(full_path), sub_rel, fn) if sub_rel != '.' else os.path.join(os.path.basename(full_path), fn)
                        zf.write(file_full, arcname)
            else:
                # Single file
                zf.write(full_path, os.path.basename(full_path))
    
    buf.seek(0)
    return send_file(buf, mimetype='application/zip', as_attachment=True, download_name='nas_download.zip')

# === SHARE LINK API ===

SHARE_DIR = os.path.join(ROOT_DIR, '.shares')

def get_share_meta():
    if not os.path.exists(os.path.join(SHARE_DIR, 'shares.json')):
        return {'shares': []}
    try:
        with open(os.path.join(SHARE_DIR, 'shares.json'), 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return {'shares': []}

def save_share_meta(meta):
    os.makedirs(SHARE_DIR, exist_ok=True)
    with open(os.path.join(SHARE_DIR, 'shares.json'), 'w', encoding='utf-8') as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

def cleanup_expired_shares(shares):
    """Remove expired shares and their physical files."""
    now = time.time()
    active = []
    for s in shares:
        if now <= s.get('expires_at', 0):
            active.append(s)
        else:
            # Remove symlink/zip if exists
            tok = s['token']
            for fname in (tok, f'{tok}.zip'):
                fpath = os.path.join(SHARE_DIR, fname)
                try:
                    if os.path.islink(fpath) or os.path.exists(fpath):
                        os.unlink(fpath)
                except Exception:
                    pass
    return active

import secrets

def generate_token():
    return secrets.token_urlsafe(12)

@nas_bp.route('/api/shares/create', methods=['POST'])
@require_auth
def shares_create():
    """Create a temporary share link for a file/folder."""
    data = request.json
    paths = data.get('paths', [])
    path = data.get('path', '')
    expires_hours = int(data.get('expires_hours', 24))
    password = data.get('password', '')
    
    # Cleanup expired shares first
    meta = get_share_meta()
    meta['shares'] = cleanup_expired_shares(meta['shares'])
    save_share_meta(meta)
    
    # Handle batch: multiple paths → create a ZIP bundle
    if len(paths) > 1:
        from io import BytesIO
        import zipfile
        os.makedirs(SHARE_DIR, exist_ok=True)
        token = generate_token()
        expires_at = time.time() + (expires_hours * 3600)
        
        # Resolve and validate all paths
        resolved = []
        for p in paths:
            fp = resolve_path(p)
            if not fp.startswith(ROOT_DIR):
                return jsonify({"error": f"Forbidden: {p}"}), 403
            if not os.path.exists(fp):
                return jsonify({"error": f"Not found: {p}"}), 404
            resolved.append((p, fp))
        
        # Create ZIP in shares dir
        zip_name = f'{token}.zip'
        zip_path = os.path.join(SHARE_DIR, zip_name)
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for orig_path, full_path in resolved:
                if os.path.isdir(full_path):
                    for dirpath, dirnames, filenames in os.walk(full_path):
                        for fn in filenames:
                            file_full = os.path.join(dirpath, fn)
                            arcname = os.path.relpath(file_full, os.path.dirname(full_path))
                            zf.write(file_full, arcname)
                else:
                    zf.write(full_path, os.path.basename(full_path))
        
        meta = get_share_meta()
        meta['shares'].append({
            'token': token,
            'original_path': f'batch_{zip_name}',
            'is_dir': False,
            'is_batch': True,
            'created_at': time.strftime('%Y-%m-%dT%H:%M:%S'),
            'expires_at': expires_at,
            'has_password': bool(password),
            'password': password
        })
        save_share_meta(meta)
        return jsonify({
            'success': True,
            'token': token,
            'expires_at': expires_at,
            'url': f'/nas/share/{token}'
        })
    
    # Single file share (also used when paths has 1 item)
    if len(paths) == 1:
        path = paths[0]
    # fall through to single path
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR):
        return jsonify({"error": "Forbidden"}), 403
    if not os.path.exists(full_path):
        return jsonify({"error": "File not found"}), 404
    
    # Generate token
    token = generate_token()
    expires_at = time.time() + (expires_hours * 3600)
    
    # Create symlink in shares dir
    os.makedirs(SHARE_DIR, exist_ok=True)
    link_path = os.path.join(SHARE_DIR, token)
    
    # Use symlink if possible, else copy
    rel_path = os.path.relpath(full_path, ROOT_DIR)
    try:
        os.symlink(full_path, link_path)
        link_type = 'symlink'
    except Exception:
        link_type = 'copy'
    
    # Save metadata
    meta = get_share_meta()
    meta['shares'].append({
        'token': token,
        'original_path': rel_path,
        'is_dir': os.path.isdir(full_path),
        'created_at': time.strftime('%Y-%m-%dT%H:%M:%S'),
        'expires_at': expires_at,
        'has_password': bool(password),
        'password': password
    })
    save_share_meta(meta)
    
    return jsonify({
        'success': True,
        'token': token,
        'expires_at': expires_at,
        'url': f'/nas/share/{token}'
    })


@nas_bp.route('/share/<token>')
def share_access(token):
    """Access a shared file/folder (no auth required)."""
    meta = get_share_meta()
    share = None
    for s in meta.get('shares', []):
        if s['token'] == token:
            share = s
            break
    
    if not share:
        return jsonify({"error": "Share link not found"}), 404
    
    # Check expiration
    if time.time() > share.get('expires_at', 0):
        return jsonify({"error": "Share link has expired"}), 410
    
    # Check password if set
    pw = request.args.get('password', '')
    error_msg = None
    if share.get('password'):
        if not pw:
            error_msg = 'This share is password-protected'
        elif pw != share['password']:
            error_msg = 'Incorrect password'
        
        if error_msg:
            err_display = '❌ Invalid password' if 'Incorrect' in error_msg else '🔐 Password required'
            share_html = f'''<!DOCTYPE html>
<html lang="zh-TW">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>J.NAS - Share Access</title>
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; background: #0a0a14; color: #e0e0ff; display:flex; align-items:center; justify-content:center; min-height:100vh; }}
.card {{ background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.06); border-radius:16px; padding:2rem; width:90%; max-width:380px; text-align:center; }}
.icon {{ font-size:48px; margin-bottom:12px; }}
h1 {{ font-size:18px; font-weight:600; margin-bottom:4px; }}
.sub {{ font-size:13px; color: rgba(255,255,255,0.35); margin-bottom:20px; word-break: break-all; }}
input {{ width:100%; background:rgba(255,255,255,0.05); border:1px solid rgba(255,255,255,0.1); color:#fff; padding:10px 14px; border-radius:10px; font-size:14px; outline:none; margin-bottom:12px; box-sizing:border-box; }}
input:focus {{ border-color:#00f2ff; }}
button {{ width:100%; background:rgba(0,242,255,0.1); border:1px solid rgba(0,242,255,0.15); color:#00f2ff; padding:10px; border-radius:10px; cursor:pointer; font-size:14px; font-weight:600; transition:background 0.2s; }}
button:hover {{ background:rgba(0,242,255,0.2); }}
.err {{ background:rgba(255,50,50,0.08); border:1px solid rgba(255,50,50,0.15); color:#ff5252; padding:8px 12px; border-radius:8px; font-size:12px; margin-bottom:12px; }}
</style>
</head>
<body>
<div class="card">
<div class="icon">{'📦' if share.get('is_batch') else '📁' if share.get('is_dir') else '📄'}</div>
'''
            if error_msg == 'Incorrect password':
                share_html += '<div class="err">❌ Incorrect password</div>'
            else:
                share_html += '<div class="err">🔐 This share is password-protected</div>'
            share_html += f'''<h1>Shared {share['original_path']}</h1>
<form method="get" action="">
<input type="password" name="password" placeholder="Enter share password" autofocus>
<button type="submit">🔓 Unlock & Download</button>
</form>
</div>
</body>
</html>'''
            return share_html, 401
    
    # Resolve the file
    if share.get('is_batch'):
        # Batch share: serve the pre-built ZIP from shares dir
        zip_path = os.path.join(SHARE_DIR, f"{share['token']}.zip")
        if not os.path.exists(zip_path):
            return jsonify({"error": "Batch share file no longer exists"}), 404
        return send_file(zip_path, mimetype='application/zip', as_attachment=True, download_name='share_bundle.zip')
    
    original_full = os.path.join(ROOT_DIR, share['original_path'])
    if not os.path.exists(original_full):
        return jsonify({"error": "Original file no longer exists"}), 404
    
    if share.get('is_dir'):
        # Serve directory as ZIP
        from io import BytesIO
        import zipfile
        buf = BytesIO()
        dirname = os.path.basename(original_full)
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
            for dirpath, dirnames, filenames in os.walk(original_full):
                for fn in filenames:
                    file_full = os.path.join(dirpath, fn)
                    arcname = os.path.relpath(file_full, os.path.dirname(original_full))
                    zf.write(file_full, arcname)
        buf.seek(0)
        return send_file(buf, mimetype='application/zip', as_attachment=True, download_name=f'{dirname}.zip')
    
    return send_file(original_full, as_attachment=True, download_name=os.path.basename(original_full))


@nas_bp.route('/api/shares/list', methods=['GET'])
@require_auth
def shares_list():
    """List active share links."""
    meta = get_share_meta()
    # Clean up expired before listing
    meta['shares'] = cleanup_expired_shares(meta['shares'])
    save_share_meta(meta)
    now = time.time()
    # Filter out expired
    active = []
    for s in meta.get('shares', []):
        if now <= s.get('expires_at', 0):
            active.append({
                'token': s['token'],
                'original_path': s['original_path'],
                'is_dir': s.get('is_dir', False),
                'created_at': s.get('created_at', ''),
                'expires_at': s['expires_at'],
                'has_password': s.get('has_password', False),
                'url': f'/nas/share/{s["token"]}'
            })
    return jsonify({'shares': active})


@nas_bp.route('/api/shares/delete', methods=['POST'])
@require_auth
def shares_delete():
    """Delete a share link."""
    token = request.json.get('token', '')
    meta = get_share_meta()
    meta['shares'] = [s for s in meta.get('shares', []) if s['token'] != token]
    save_share_meta(meta)
    
    # Remove symlink if exists
    link_path = os.path.join(SHARE_DIR, token)
    if os.path.exists(link_path):
        try:
            if os.path.islink(link_path):
                os.unlink(link_path)
            elif os.path.isfile(link_path):
                os.remove(link_path)
        except Exception:
            pass
    
    return jsonify({'success': True})


# === WEBDAV ===

WEBDAV_PORT = 8001
webdav_process = None

def start_webdav():
    """Start wsgidav server as a subprocess."""
    global webdav_process
    if webdav_process and webdav_process.poll() is None:
        return True  # Already running
    
    import subprocess
    venv_python = sys.executable
    webdav_script = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'webdav_server.py')
    
    # Ensure the webdav server script exists and is executable
    if not os.path.exists(webdav_script):
        logger.error(f"WebDAV script not found at {webdav_script}")
        return False
    os.chmod(webdav_script, 0o755)
    
    try:
        webdav_process = subprocess.Popen(
            [venv_python, webdav_script],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
        return True
    except Exception as e:
        logger.error(f"Failed to start WebDAV: {e}")
        return False

def stop_webdav():
    """Stop the WebDAV subprocess."""
    global webdav_process
    if webdav_process and webdav_process.poll() is None:
        webdav_process.terminate()
        try:
            webdav_process.wait(timeout=5)
        except subprocess.TimeoutExpired:
            webdav_process.kill()
        webdav_process = None
        return True
    return False


def is_webdav_running():
    """Check if WebDAV subprocess is running."""
    global webdav_process
    return webdav_process is not None and webdav_process.poll() is None


@nas_bp.route('/api/webdav/status', methods=['GET'])
@require_auth
def webdav_status():
    """Get WebDAV server status."""
    return jsonify({
        'running': is_webdav_running(),
        'port': WEBDAV_PORT,
        'url': f'http://100.99.4.98:{WEBDAV_PORT}/',
        'dav_url': f'http://100.99.4.98:{WEBDAV_PORT}/',
        'username': 'jerry'
    })


@nas_bp.route('/api/webdav/start', methods=['POST'])
@require_auth
def webdav_start():
    """Start WebDAV server."""
    if is_webdav_running():
        return jsonify({'success': True, 'message': 'Already running'})
    ok = start_webdav()
    return jsonify({'success': ok, 'running': ok})


@nas_bp.route('/api/webdav/stop', methods=['POST'])
@require_auth
def webdav_stop():
    """Stop WebDAV server."""
    stop_webdav()
    return jsonify({'success': True, 'running': False})


@nas_bp.route('/api/shutdown', methods=['POST'])
@require_auth
def shutdown_server():
    logging.getLogger(__name__).info('Shutdown request received. Shutting down server...')
    
    # Use Werkzeug's built-in shutdown mechanism instead of os._exit(0)
    # which kills the entire Python process (including the GUI).
    func = request.environ.get('werkzeug.server.shutdown')
    if func:
        func()
        return jsonify({"success": True, "message": "Server shutting down..."})
    else:
        # Fallback: use os._exit only if Werkzeug shutdown is unavailable
        def exit_process():
            time.sleep(1)
            os._exit(0)
        threading.Thread(target=exit_process).start()
        return jsonify({"success": True, "message": "Server shutting down (emergency)..."})

@nas_bp.route('/api/view')
@require_auth
def view_file():
    path = request.args.get('path', '')
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR): return jsonify({"error": "Forbidden"}), 403
    return send_from_directory(os.path.dirname(full_path), os.path.basename(full_path), as_attachment=False)

# ====== DOCUMENT READ/WRITE APIs ======

@nas_bp.route('/api/read_doc')
@require_auth
def read_document():
    """Read Office/PDF document, return structured JSON for frontend card."""
    path = request.args.get('path', '')
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR):
        return jsonify({"error": "Forbidden"}), 403

    ext = os.path.splitext(full_path)[1].lower()

    try:
        if ext == '.docx':
            from docx import Document
            doc = Document(full_path)
            paragraphs = []
            for p in doc.paragraphs:
                if p.text.strip():
                    paragraphs.append({
                        "text": p.text,
                        "style": p.style.name if p.style else "Normal"
                    })
            return jsonify({
                "type": "word",
                "content": paragraphs,
                "filepath": path,
                "filename": os.path.basename(full_path)
            })

        elif ext in ('.xls', '.xlsx'):
            import openpyxl
            wb = openpyxl.load_workbook(full_path, data_only=True)
            sheets = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                data = []
                for row in ws.iter_rows(values_only=True):
                    data.append([str(c) if c is not None else "" for c in row])
                sheets[sheet_name] = data
            wb.close()
            return jsonify({
                "type": "excel",
                "sheets": sheets,
                "activeSheet": wb.sheetnames[0] if wb.sheetnames else "",
                "sheetNames": wb.sheetnames,
                "filepath": path,
                "filename": os.path.basename(full_path)
            })

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(full_path)
            slides = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            t = p.text.strip()
                            if t:
                                texts.append(t)
                slides.append({"texts": texts})
            return jsonify({
                "type": "powerpoint",
                "slides": slides,
                "slideCount": len(slides),
                "filepath": path,
                "filename": os.path.basename(full_path)
            })

        elif ext == '.pdf':
            import pypdf
            reader = pypdf.PdfReader(full_path)
            pages = []
            for page in reader.pages:
                pages.append(page.extract_text())
            return jsonify({
                "type": "pdf",
                "pages": pages,
                "pageCount": len(pages),
                "filepath": path,
                "filename": os.path.basename(full_path)
            })

        return jsonify({"error": f"Unsupported: {ext}"}), 400

    except Exception as e:
        logger.error(f"read_doc error: {e}")
        return jsonify({"error": str(e)}), 500


@nas_bp.route('/api/pdf_page')
@require_auth
def pdf_page_image():
    """Render one PDF page as JPEG image using pypdfium2."""
    path = request.args.get('path', '')
    page_num = int(request.args.get('page', 0))
    scale = float(request.args.get('scale', '1.5'))

    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR):
        return jsonify({"error": "Forbidden"}), 403

    try:
        import pypdfium2 as pdfium
        pdf = pdfium.PdfDocument(full_path)
        if page_num >= len(pdf):
            return jsonify({"error": "Page not found"}), 404

        page = pdf[page_num]
        bitmap = page.render(scale=scale)
        pil_image = bitmap.to_pil()
        img_io = io.BytesIO()
        pil_image.save(img_io, 'JPEG', quality=92)
        img_io.seek(0)
        return send_file(img_io, mimetype='image/jpeg')
    except Exception as e:
        logger.error(f"pdf_page error: {e}")
        return jsonify({"error": str(e)}), 500


@nas_bp.route('/api/save_doc', methods=['POST'])
@require_auth
def save_document():
    """Save document from JSON payload (round-trip editing)."""
    data = request.json
    path = data.get('path', '')
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR):
        return jsonify({"error": "Forbidden"}), 403

    ext = os.path.splitext(full_path)[1].lower()

    try:
        if ext == '.docx':
            from docx import Document
            doc = Document()
            for para in data.get('content', []):
                if isinstance(para, dict):
                    doc.add_paragraph(para.get('text', ''))
                else:
                    doc.add_paragraph(str(para))
            doc.save(full_path)
            return jsonify({"success": True})

        elif ext in ('.xls', '.xlsx'):
            import openpyxl
            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            sheets_data = data.get('sheets', {})
            for sheet_name, rows in sheets_data.items():
                ws = wb.create_sheet(title=sheet_name[:31])  # Excel sheet name limit
                for row in rows:
                    ws.append(row)
            wb.save(full_path)
            return jsonify({"success": True})

        elif ext == '.pptx':
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation()
            for slide_data in data.get('slides', []):
                slide_layout = prs.slide_layouts[6]  # blank
                slide = prs.slides.add_slide(slide_layout)
                for text in slide_data.get('texts', []):
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
                    tf = txBox.text_frame
                    tf.text = text
            prs.save(full_path)
            return jsonify({"success": True})

        return jsonify({"error": f"Unsupported: {ext}"}), 400

    except Exception as e:
        logger.error(f"save_doc error: {e}")
        return jsonify({"error": str(e)}), 500


@nas_bp.route('/api/diag_ffmpeg')
@require_auth
def diag_ffmpeg():
    """Diagnostic: show ffmpeg resolution state and run smoke test."""
    import sys, os, shutil, subprocess
    
    # Clear cache so we can test fresh
    global _FFMPEG_CACHED
    _FFMPEG_CACHED = None

    which_fp = shutil.which('ffmpeg')
    bundled_fp = None
    if getattr(sys, 'frozen', False):
        bundled_fp = os.path.join(sys._MEIPASS, 'nas', 'bin', 'ffmpeg.exe')

    # Run the full resolution flow
    resolved, si = _resolve_ffmpeg()

    # Smoke test the resolved binary
    smoke = None
    if resolved:
        try:
            r = subprocess.run([resolved, '-version'], capture_output=True, timeout=5, startupinfo=si)
            version_out = r.stdout.decode('utf-8', errors='replace').split('\n')[0] if r.stdout else ''
            smoke = {'ok': r.returncode == 0, 'version': version_out, 'exit_code': r.returncode}
        except Exception as e:
            smoke = {'ok': False, 'error': str(e)}

    info = {
        'sys.frozen': getattr(sys, 'frozen', False),
        'sys._MEIPASS': getattr(sys, '_MEIPASS', None),
        'sys.platform': sys.platform,
        'shutil.which(ffmpeg)': which_fp,
        'bundled_path': bundled_fp,
        'bundled_exists': os.path.exists(bundled_fp) if bundled_fp else None,
        'resolved_path': resolved,
        'smoke_test': smoke,
    }
    return jsonify(info)


_FFMPEG_CACHED = None

def _resolve_ffmpeg():
    """Resolve & verify ffmpeg: system PATH first (winget), then raw cmd.
    On Windows, returns startupinfo to suppress the console window flash.
    Result cached after first successful verification."""
    global _FFMPEG_CACHED
    if _FFMPEG_CACHED is not None:
        return _FFMPEG_CACHED

    startupinfo = None
    if sys.platform == 'win32':
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags = subprocess.STARTF_USESHOWWINDOW
        startupinfo.wShowWindow = 0

    # Prefer system PATH (winget-installed) over bundled binary
    ffmpeg_bin = shutil.which('ffmpeg')
    if ffmpeg_bin:
        _FFMPEG_CACHED = (ffmpeg_bin, startupinfo)
        return _FFMPEG_CACHED

    # Fallback: bundled binary in _MEIPASS
    if getattr(sys, 'frozen', False):
        bundled = os.path.join(sys._MEIPASS, 'nas', 'bin', 'ffmpeg.exe')
        if os.path.exists(bundled) and _verify_ffmpeg(bundled, startupinfo):
            _FFMPEG_CACHED = (bundled, startupinfo)
            return _FFMPEG_CACHED

    # Linux brew fallback
    if os.path.exists('/home/linuxbrew/.linuxbrew/bin/ffmpeg'):
        ffmpeg_bin = '/home/linuxbrew/.linuxbrew/bin/ffmpeg'
    else:
        ffmpeg_bin = 'ffmpeg'

    _FFMPEG_CACHED = (ffmpeg_bin, startupinfo)
    return _FFMPEG_CACHED


def _verify_ffmpeg(path, startupinfo):
    """Quick smoke-test: run 'ffmpeg -version' to confirm binary works."""
    try:
        r = subprocess.run([path, '-version'], capture_output=True, timeout=5, startupinfo=startupinfo)
        return r.returncode == 0
    except Exception:
        return False


@nas_bp.route('/api/thumbnail')
def get_thumbnail():
    path = request.args.get('path', '')
    full_path = resolve_path(path)
    if not full_path.startswith(ROOT_DIR):
        return redirect('https://cdn-icons-png.flaticon.com/512/2961222.png', code=302)
    
    ext = os.path.splitext(full_path)[1].lower()
    
    try:
        # 1. Image Thumbnails (Pure Memory)
        if ext in IMG_EXTS:
            from PIL import Image, ImageOps
            img = Image.open(full_path)
            img = ImageOps.exif_transpose(img)
            img.thumbnail((800, 800), Image.LANCZOS)
            if img.mode in ('RGBA', 'LA', 'P'):
                bg = Image.new('RGB', img.size, (30, 30, 40))
                if img.mode == 'RGBA':
                    bg.paste(img, mask=img.split()[3])
                else:
                    bg.paste(img)
                img = bg
            
            img_io = io.BytesIO()
            img.save(img_io, 'JPEG', quality=90)
            img_io.seek(0)
            return send_file(img_io, mimetype='image/jpeg')

        # 2. Video Thumbnails (Pure Memory via Pipe)
        if ext in VID_EXTS:
            import shutil
            ffmpeg_bin, si = _resolve_ffmpeg()
            for timestamp in ['00:00:01', '00:00:00']:
                cmd = [
                    ffmpeg_bin, '-loglevel', 'error', '-ss', timestamp, 
                    '-i', full_path, '-vframes', '1', '-f', 'image2pipe', 
                    '-vcodec', 'mjpeg', '-vf', 'scale=600:-1', '-'
                ]
                try:
                    result = subprocess.run(cmd, capture_output=True, timeout=10,
                                            startupinfo=si)
                    if result.returncode == 0 and result.stdout:
                        return Response(result.stdout, mimetype='image/jpeg')
                except Exception as ve:
                    err_detail = str(ve)
                    if result and hasattr(result, 'stderr') and result.stderr:
                        err_detail += ' | STDERR: ' + result.stderr.decode('utf-8', errors='replace')
                    return jsonify({"error": f"ffmpeg failed at {timestamp}", "detail": err_detail}), 500

    except Exception as e:
        logger.error(f"General thumb error: {e}")
        pass

    # 3. Special Document Icons
    icon_paths = {
        '.doc': ('nas_tool/nas/icons/word---67ab8593-d97f-4d34-bea6-25bc16087a34.png', 'https://cdn-icons-png.flaticon.com/512/732220.png'),
        '.docx': ('nas_tool/nas/icons/word---67ab8593-d97f-4d34-bea6-25bc16087a34.png', 'https://cdn-icons-png.flaticon.com/512/732220.png'),
        '.xls': ('nas_tool/nas/icons/logo---fbd73b10-ff25-4e86-b0f0-ac7a0ae3d93c.png', 'https://cdn-icons-png.flaticon.com/512/732222.png'),
        '.xlsx': ('nas_tool/nas/icons/logo---fbd73b10-ff25-4e86-b0f0-ac7a0ae3d93c.png', 'https://cdn-icons-png.flaticon.com/512/732222.png'),
        '.ppt': ('nas_tool/nas/icons/logo---0a771b7b-b41e-4235-bf90-003e8c46e5a6.png', 'https://cdn-icons-png.flaticon.com/512/732225.png'),
        '.pptx': ('nas_tool/nas/icons/logo---0a771b7b-b41e-4235-bf90-003e8c46e5a6.png', 'https://cdn-icons-png.flaticon.com/512/732225.png'),
    }
    if ext in icon_paths:
        custom_path, fallback_url = icon_paths[ext]
        full_custom = os.path.join(ROOT_DIR, custom_path)
        if os.path.exists(full_custom):
            return send_file(full_custom, mimetype='image/png')
        return redirect(fallback_url, code=302)

    # 4. Generic Fallback Icons
    icon_map = {
        'jpg':'337943','jpeg':'337943','png':'337943','gif':'337943','webp':'337943','svg':'337943',
        'mp4':'1179067','mov':'1179067','avi':'1179067','mkv':'1179067','webm':'1179067',
        'mp3':'461261','wav':'461261','m4a':'461261','aac':'461261','flac':'461261',
        'pdf':'337946',
        'doc':'https://cdn-icons-png.flaticon.com/512/732/732220.png',
        'docx':'https://cdn-icons-png.flaticon.com/512/732/732220.png',
        'xls':'https://cdn-icons-png.flaticon.com/512/732/732222.png',
        'xlsx':'https://cdn-icons-png.flaticon.com/512/732/732222.png',
        'ppt':'https://cdn-icons-png.flaticon.com/512/732/732225.png',
        'pptx':'https://cdn-icons-png.flaticon.com/512/732/732225.png',
        'py':'1055644','js':'1055644','html':'1055644','css':'1055644','json':'1055644',
        'md':'1055644','txt':'1055644','log':'1055644','sh':'1055644',
        'zip':'2961218','tar':'2961218','gz':'2961218','rar':'2961218'
    }
    
    icon_val = icon_map.get(ext[1:] if ext.startswith('.') else ext, '2961222')
    if icon_val.startswith('http'):
        return redirect(icon_val, code=302)
    return redirect(f'https://cdn-icons-png.flaticon.com/512/{icon_val}.png', code=302)

@nas_bp.route('/manifest.json')
def pwa_manifest():
    return send_from_directory(os.path.dirname(os.path.abspath(__file__)), 'manifest.json')

@nas_bp.route('/sw.js')
def pwa_sw():
    return send_from_directory(os.path.dirname(os.path.abspath(__file__)), 'sw.js')

app.register_blueprint(nas_bp)

@app.route('/', defaults={'path': ''})
@app.route('/<path:path>')
def proxy_to_finance(path):
    # Ensure NAS routes are handled by the blueprint
    if path == 'nas' or path.startswith('nas/'):
        return redirect('/nas/', code=301)
    
    finance_url = f'http://localhost:5000/{path}'
    try:
        # Proxy the request to the Finance server
        resp = requests.request(
            method=request.method,
            url=finance_url,
            headers={k: v for k, v in request.headers if k.lower() != 'host'},
            data=request.get_data(),
            cookies=request.cookies,
            allow_redirects=False
        )
        
        # Filter out hop-by-hop headers and headers the proxy will regenerate
        excluded_headers = ['content-encoding', 'content-length', 'transfer-encoding', 'connection', 'date', 'server']
        headers = [(k, v) for k, v in resp.raw.headers.items() if k.lower() not in excluded_headers]
        
        return Response(resp.content, resp.status_code, headers)
    except Exception as e:
        logger.error(f"Proxy error: {e}")
        return jsonify({"error": "Finance server unreachable"}), 502

if __name__ == '__main__':
    # When run directly (e.g. via jnexus.service), go through init_app()
    # so the root path is handled consistently. Override via env vars.
    root = os.environ.get('NAS_ROOT_DIR', '/home/jerry/workspace')
    password = os.environ.get('NAS_PASSWORD', 'JERRY_NEXUS_2026')
    port = int(os.environ.get('NAS_PORT', '8000'))
    init_app(root, password, port)
    app.run(host='0.0.0.0', port=port)
